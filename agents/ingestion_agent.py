

# agents/ingestion_agent.py

import fitz  # PyMuPDF
import docx
import pandas as pd
import os
from pptx import Presentation

from utils.mcp import parse_message, create_message

class IngestionAgent:
    def __init__(self, name="IngestionAgent"):
        self.name = name

    def extract_text_from_file(self, file_path):
        ext = file_path.split('.')[-1]
        if ext == "pdf":
            return self.extract_text_from_pdf(file_path)
        elif ext == "docx":
            return self.extract_text_from_docx(file_path)
        elif ext == "pptx":
            return self.extract_text_from_pptx(file_path)
        elif ext == "csv":
            return self.extract_text_from_csv(file_path)
        elif ext in ["txt", "md"]:
            return self.extract_text_from_txt(file_path)
        else:
            return ""

    def extract_text_from_pdf(self, file_path):
        text = ""
        with fitz.open(file_path) as doc:
            for page in doc:
                text += page.get_text()
        return text

    def extract_text_from_docx(self, file_path):
        doc = docx.Document(file_path)
        return "\n".join([p.text for p in doc.paragraphs])

    def extract_text_from_pptx(self, file_path):
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    def extract_text_from_csv(self, file_path):
        df = pd.read_csv(file_path)
        return df.to_string(index=False)

    def extract_text_from_txt(self, file_path):
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()

    def chunk_text(self, text, chunk_size=300, overlap=50):
        words = text.split()
        chunks = []
        for i in range(0, len(words), chunk_size - overlap):
            chunk = " ".join(words[i:i+chunk_size])
            chunks.append(chunk)
        return chunks

    def handle_message(self, message):
        sender, receiver, msg_type, trace_id, payload = parse_message(message)
        file_path = payload.get("file_path")

        if not os.path.exists(file_path):
            return create_message(
                self.name, sender, "error", trace_id,
                {"error": f"File not found: {file_path}"}
            )

        text = self.extract_text_from_file(file_path)
        chunks = self.chunk_text(text)

        return create_message(
            self.name, sender, "ingested", trace_id,
            {"chunks": chunks}
        )

