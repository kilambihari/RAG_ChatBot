import os
import fitz             # PyMuPDF
import pandas as pd
from docx import Document
from pptx import Presentation
import markdown         # only if you really want to convert md → html later

def parse_pdf(file_path: str, chunk_size: int = 1000, chunk_overlap: int = 200) -> list[str]:
    """Extract text from PDF using PyMuPDF (fitz) - much faster & more reliable than PyPDF2"""
    try:
        doc = fitz.open(file_path)
        full_text = ""
        for page in doc:
            text = page.get_text("text").strip()
            if text:
                full_text += text + "\n\n"
        doc.close()

        if not full_text.strip():
            return []

        # Simple overlapping chunking
        chunks = []
        start = 0
        while start < len(full_text):
            end = start + chunk_size
            chunks.append(full_text[start:end])
            start = end - chunk_overlap
        return chunks

    except Exception as e:
        raise ValueError(f"Failed to parse PDF {file_path}: {str(e)}")


def parse_docx(file_path: str) -> list[str]:
    """Extract paragraphs from .docx"""
    try:
        doc = Document(file_path)
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        # You could also extract tables here if needed
        return paragraphs
    except Exception as e:
        raise ValueError(f"Failed to parse DOCX {file_path}: {str(e)}")


def parse_txt_or_md(file_path: str, is_markdown: bool = False) -> list[str]:
    """Handle .txt and .md files uniformly"""
    try:
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            content = f.read()

        if is_markdown:
            # Optional: convert markdown to plain text if you want
            # html = markdown.markdown(content)
            # plain = ... (use html2text or just keep raw)
            pass

        # Split into non-empty lines or paragraphs
        lines = [line.strip() for line in content.splitlines() if line.strip()]
        return lines

    except Exception as e:
        raise ValueError(f"Failed to parse text/markdown file {file_path}: {str(e)}")


def parse_csv(file_path: str) -> list[str]:
    """Convert CSV rows to readable strings"""
    try:
        df = pd.read_csv(file_path, dtype=str, keep_default_na=False)
        # Join columns with | separator
        rows = [" | ".join(row.astype(str).tolist()) for _, row in df.iterrows()]
        return [row for row in rows if row.strip()]
    except Exception as e:
        raise ValueError(f"Failed to parse CSV {file_path}: {str(e)}")


def parse_pptx(file_path: str) -> list[str]:
    """Extract text from PowerPoint slides"""
    try:
        prs = Presentation(file_path)
        slide_texts = []

        for slide in prs.slides:
            slide_content = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text.strip())
            if slide_content:
                slide_texts.append(" ".join(slide_content))

        return slide_texts

    except Exception as e:
        raise ValueError(f"Failed to parse PPTX {file_path}: {str(e)}")


def parse_document(file_path: str, chunk_size: int = 1000, chunk_overlap: int = 200) -> list[str]:
    """
    Main entry point - returns list of text chunks/paragraphs depending on file type
    """
    ext = os.path.splitext(file_path)[1].lower()

    parsers = {
        ".pdf":  lambda: parse_pdf(file_path, chunk_size, chunk_overlap),
        ".docx": lambda: parse_docx(file_path),
        ".txt":  lambda: parse_txt_or_md(file_path, is_markdown=False),
        ".md":   lambda: parse_txt_or_md(file_path, is_markdown=True),
        ".csv":  lambda: parse_csv(file_path),
        ".pptx": lambda: parse_pptx(file_path),
    }

    if ext not in parsers:
        raise ValueError(f"Unsupported file extension: {ext}")

    return parsers[ext]()


# Optional: small test helper (can be removed later)
if __name__ == "__main__":
    # Example usage
    try:
        chunks = parse_document("example.pdf")
        print(f"Extracted {len(chunks)} chunks")
        print(chunks[0][:200] if chunks else "No content")
    except Exception as e:
        print(e)

